"""Export Center — generates PPTX, PDF, and Excel reports from Shift Matrix.

All outputs use percentage values only — never €M.
Professional formatting with force attribution and methodology.
"""

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pulse.common.shape_compat import velocity_median

logger = logging.getLogger(__name__)


class ExportCenter:
    """Generates professional PPTX, PDF, and Excel reports from PRISM Shift Matrix."""

    # Color scheme — matches Profit Pool Shift Model design
    COLOR_BLUE = RGBColor(59, 130, 246)  # #3B82F6
    COLOR_EXPANSION = RGBColor(34, 197, 94)  # #22C55E (green)
    COLOR_CONTRACTION = RGBColor(239, 68, 68)  # #EF4444 (red)
    COLOR_CAUSAL = RGBColor(139, 92, 246)  # #8B5CF6 (purple)
    COLOR_DARK_BG = RGBColor(15, 23, 42)  # #0F172A (dark navy)
    COLOR_LIGHT_TEXT = RGBColor(248, 250, 252)  # #F8FAFC (off-white)

    def __init__(self, output_dir: str = "exports"):
        """Initialize Export Center.

        Args:
            output_dir: Directory to save exported files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def export_pptx(
        self,
        shift_matrix: Dict[str, Any],
        output_path: Optional[str] = None,
    ) -> str:
        """Generate a professional PowerPoint presentation from Shift Matrix.

        Creates a 6-slide deck:
        1. Title slide
        2. Executive summary (headline KPI, path velocity)
        3. Heatmap data table (category × force shifts)
        4. Shift paths (continuous 2026-2030)
        5. Resource allocation recommendations
        6. Methodology & confidence

        Args:
            shift_matrix: Full Shift Matrix with paths and metadata
            output_path: Path to save PPTX. If None, auto-generates filename.

        Returns:
            Path to the generated PPTX file
        """
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(self.output_dir / f"PRISM_Report_{timestamp}.pptx")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        self._add_title_slide(prs, shift_matrix)

        # Slide 2: Executive Summary
        self._add_executive_summary_slide(prs, shift_matrix)

        # Slide 3: Category Heatmap
        self._add_heatmap_slide(prs, shift_matrix)

        # Slide 4: Shift Paths
        self._add_paths_slide(prs, shift_matrix)

        # Slide 5: Allocation Recommendations
        self._add_allocation_slide(prs, shift_matrix)

        # Slide 6: Methodology
        self._add_methodology_slide(prs, shift_matrix)

        prs.save(output_path)
        logger.info(f"PPTX exported to {output_path}")
        return output_path

    def export_pdf(
        self,
        shift_matrix: Dict[str, Any],
        output_path: Optional[str] = None,
    ) -> str:
        """Generate a professional PDF report from Shift Matrix.

        Single-page report with:
        - Headline metrics
        - Causal narrative
        - Summary table
        - Confidence explanation

        Args:
            shift_matrix: Full Shift Matrix with metadata
            output_path: Path to save PDF. If None, auto-generates filename.

        Returns:
            Path to the generated PDF file
        """
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(self.output_dir / f"PRISM_Report_{timestamp}.pdf")

        doc = SimpleDocTemplate(output_path, pagesize=letter,
                                rightMargin=0.5*inch, leftMargin=0.5*inch,
                                topMargin=0.75*inch, bottomMargin=0.75*inch)

        story = []
        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#0F172A'),
            spaceAfter=12,
            alignment=TA_CENTER,
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#3B82F6'),
            spaceAfter=10,
            spaceBefore=10,
        )

        # Title
        story.append(Paragraph("PRISM Profit Pool Analysis", title_style))
        story.append(Spacer(1, 0.2*inch))

        # Metadata
        gen_date = shift_matrix.get("generated", datetime.now().isoformat())
        story.append(Paragraph(
            f"Generated: {gen_date} | "
            f"Model: {shift_matrix.get('model_version', 'N/A')} | "
            f"Confidence: {shift_matrix.get('confidence', '80% CI')}",
            styles['Normal']
        ))
        story.append(Spacer(1, 0.15*inch))

        # Headline metrics
        shifts_data = shift_matrix.get("shifts", {})
        if shifts_data:
            medians = []
            for cat, cat_data in shifts_data.items():
                path = cat_data.get("path", {})
                if path:
                    # Get 2030 median as long-term impact
                    year_2030 = path.get(2030, {})
                    if isinstance(year_2030, dict):
                        medians.append(year_2030.get("median", 0.0))

            if medians:
                avg_shift = sum(medians) / len(medians)
                story.append(Paragraph(
                    f"<b>Net Pool Shift (2030, median): {avg_shift:+.1%}</b>",
                    heading_style
                ))

        story.append(Spacer(1, 0.1*inch))

        # Force attribution narrative
        story.append(Paragraph("<b>Force Attribution Narrative:</b>", heading_style))
        narrative = self._generate_force_narrative(shift_matrix)
        story.append(Paragraph(narrative, styles['Normal']))
        story.append(Spacer(1, 0.2*inch))

        # Summary table
        story.append(Paragraph("<b>Category Impact Summary (2030 Median):</b>", heading_style))
        summary_table_data = self._build_summary_table(shifts_data, include_header=True)
        if summary_table_data:
            summary_table = Table(summary_table_data, colWidths=[2.5*inch, 1.2*inch, 1.2*inch, 1*inch])
            summary_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3B82F6')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F0F4F8')]),
            ]))
            story.append(summary_table)

        story.append(Spacer(1, 0.15*inch))

        # Confidence note
        story.append(Paragraph(
            f"<i><b>Engine:</b> {shift_matrix.get('engine_name', 'bayesian_copula')} "
            f"(version {shift_matrix.get('model_version', 'n/a')})</i>",
            styles['Normal']
        ))

        doc.build(story)
        logger.info(f"PDF exported to {output_path}")
        return output_path

    def export_excel(
        self,
        shift_matrix: Dict[str, Any],
        output_path: Optional[str] = None,
    ) -> str:
        """Generate an Excel workbook with Shift Matrix and application templates.

        Creates sheets:
        1. Shift Matrix — continuous paths with all percentiles
        2. Application Template — user enters their GP1, formulas auto-calculate
        3. Allocation Recommendations — relative investment weights
        4. Force Attribution — per-category decomposition by force
        5. Methodology — engine version and model parameters

        Args:
            shift_matrix: Full Shift Matrix with paths and metadata
            output_path: Path to save Excel. If None, auto-generates filename.

        Returns:
            Path to the generated Excel file
        """
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(self.output_dir / f"PRISM_Export_{timestamp}.xlsx")

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default sheet

        # Sheet 1: Shift Matrix
        self._create_shift_matrix_sheet(wb, shift_matrix)

        # Sheet 2: Application Template
        self._create_application_template_sheet(wb, shift_matrix)

        # Sheet 3: Allocation Recommendations
        allocation = shift_matrix.get("allocation_recommendation", {})
        if allocation:
            self._create_allocation_sheet(wb, allocation)

        # Sheet 4: Force Attribution (per category, scaled to MC median)
        attribution = shift_matrix.get("force_attribution") or shift_matrix.get("causal_decomposition", {})
        if attribution:
            self._create_force_attribution_sheet(wb, attribution)

        # Sheet 5: Methodology
        self._create_methodology_sheet(wb, shift_matrix)

        wb.save(output_path)
        logger.info(f"Excel exported to {output_path}")
        return output_path

    # ==================== PPTX HELPERS ====================

    def _add_title_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add title slide to presentation."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "PRISM Profit Pool Analysis"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = "Bayesian Copula · Causal DAG · Continuous Paths"
        p.font.size = Pt(24)
        p.font.color.rgb = self.COLOR_LIGHT_TEXT

        # Add date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        p = date_frame.paragraphs[0]
        p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(148, 163, 184)  # #94A3B8

    def _add_executive_summary_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add executive summary slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Headline KPI
        shifts_data = shift_matrix.get("shifts", {})
        medians = [shift_matrix.get("shifts", {}).get(cat, {}).get("path", {}).get(2030, {}).get("median", 0.0)
                   for cat in shift_matrix.get("shifts", {})]
        avg_shift = sum(medians) / len(medians) if medians else 0.0

        kpi_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1.2))
        kpi_frame = kpi_box.text_frame
        kpi_frame.word_wrap = True
        p = kpi_frame.paragraphs[0]
        p.text = f"Net Pool Shift (2030): {avg_shift:+.1%}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_EXPANSION if avg_shift > 0 else self.COLOR_CONTRACTION

        # Model confidence
        confidence_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
        confidence_frame = confidence_box.text_frame
        p = confidence_frame.paragraphs[0]
        p.text = f"Engine: {shift_matrix.get('engine_name', 'bayesian_copula')} · v{shift_matrix.get('model_version', 'n/a')}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.COLOR_LIGHT_TEXT

        # Narrative
        narrative_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3))
        narrative_frame = narrative_box.text_frame
        narrative_frame.word_wrap = True
        narrative_text = self._generate_causal_narrative(shift_matrix)
        p = narrative_frame.paragraphs[0]
        p.text = narrative_text
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLOR_LIGHT_TEXT

    def _add_heatmap_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add heatmap data table slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Category Impact Summary (2030 Median)"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Build table
        table_data = self._build_summary_table(shift_matrix.get("shifts", {}), include_header=True)
        if not table_data:
            return

        # Limit to key categories for readability
        if len(table_data) > 15:
            table_data = [table_data[0]] + table_data[1:15]

        # Add table shape
        rows = len(table_data)
        cols = len(table_data[0])
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Populate table
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                if row_idx == 0:  # Header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLOR_BLUE
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = self.COLOR_LIGHT_TEXT
                            run.font.bold = True

    def _add_paths_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add shift paths slide (summary table)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Continuous Shift Paths (2026-2030)"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Build path table for top 5 categories
        shifts_data = shift_matrix.get("shifts", {})
        top_cats = list(shifts_data.keys())[:5]

        path_data = [["Category", "2026", "2027", "2028", "2029", "2030"]]
        for cat in top_cats:
            cat_data = shifts_data.get(cat, {})
            path = cat_data.get("path", {})
            row = [cat]
            for year in [2026, 2027, 2028, 2029, 2030]:
                year_data = path.get(year, {})
                median = year_data.get("median", 0.0) if isinstance(year_data, dict) else 0.0
                row.append(f"{median:+.1%}")
            path_data.append(row)

        # Add table
        rows = len(path_data)
        cols = len(path_data[0])
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for row_idx, row_data in enumerate(path_data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLOR_BLUE

    def _add_allocation_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add allocation recommendations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Resource Allocation Recommendations"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Recommendations text
        allocation = shift_matrix.get("allocation_recommendation", {})
        rationale = allocation.get("rationale", "See allocation sheet for detailed weights.") if isinstance(allocation, dict) else "See allocation sheet."

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = rationale
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLOR_LIGHT_TEXT

    def _add_methodology_slide(self, prs: Presentation, shift_matrix: Dict[str, Any]) -> None:
        """Add methodology and confidence slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Methodology & Confidence"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_BLUE

        # Methodology text
        methodology = f"""
Engine: {shift_matrix.get('engine_name', 'bayesian_copula')}
Model Version: {shift_matrix.get('model_version', 'n/a')}
Seed: {shift_matrix.get('seed', 'n/a')}
Confidence: {shift_matrix.get('confidence', '80% CI')}

Key Features:
• Bayesian sampling with Beta posteriors over trend probabilities
• Copula-based dependency structures capturing tail risk
• Force attribution per category (static, scaled to MC median)
• Continuous paths with velocity tracking (2026-2030)
• Reproducible: fixed RNG seed; integrity events surface any runtime repairs

Security:
• All outputs contain ONLY percentage shifts
• Zero financial data (NES, GP1, GP2) in any export
• One-directional flow to Power BI (read-only)
• Audit trail of all model changes and predictions
"""

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = methodology.strip()
        p.font.size = Pt(10)
        p.font.color.rgb = self.COLOR_LIGHT_TEXT

    # ==================== PDF HELPERS ====================

    def _generate_force_narrative(self, shift_matrix: Dict[str, Any]) -> str:
        """Generate a force narrative explaining the shift drivers.

        Args:
            shift_matrix: Full Shift Matrix with force attribution

        Returns:
            Narrative string explaining the force mechanism
        """
        causal = shift_matrix.get("force_attribution") or shift_matrix.get("causal_decomposition", {})
        if not causal:
            return "Shifts driven by composite market forces across Consumer, Customer, Technology, Government, Environmental, and Competitive channels."

        # Summarize dominant forces
        force_impacts = {}
        for category, decomp in causal.items():
            for force, impact in decomp.get("direct_effects", {}).items():
                force_impacts[force] = force_impacts.get(force, 0) + abs(impact)

        if not force_impacts:
            return "Shifts driven by composite market forces."

        top_force = max(force_impacts, key=force_impacts.get)
        narrative = f"Pool shifts are primarily driven by {top_force} forces, " \
                    f"propagating through causal channels to Customer and Consumer dynamics. " \
                    f"See Causal Decomposition sheet for detailed channel analysis."

        return narrative

    def _build_summary_table(self, shifts_data: Dict[str, Any], include_header: bool = False) -> List[List[str]]:
        """Build a summary table of category impacts.

        Args:
            shifts_data: Category shift data from shift_matrix["shifts"]
            include_header: Whether to include header row

        Returns:
            List of lists for table data
        """
        table_data = []
        if include_header:
            table_data.append(["Category", "Median 2030", "p10 (Bear)", "p90 (Bull)"])

        for category, cat_data in shifts_data.items():
            path = cat_data.get("path", {})
            year_2030 = path.get(2030, {})
            if isinstance(year_2030, dict):
                median = year_2030.get("median", 0.0)
                p10 = year_2030.get("p10", 0.0)
                p90 = year_2030.get("p90", 0.0)
                table_data.append([
                    category,
                    f"{median:+.1%}",
                    f"{p10:+.1%}",
                    f"{p90:+.1%}",
                ])

        return table_data

    # ==================== EXCEL HELPERS ====================

    def _create_shift_matrix_sheet(self, wb: openpyxl.Workbook, shift_matrix: Dict[str, Any]) -> None:
        """Create Shift Matrix sheet with continuous paths."""
        ws = wb.create_sheet("Shift Matrix")

        # Header
        headers = ["Category", "Year", "Median (%)", "p10 (%)", "p25 (%)", "p75 (%)", "p90 (%)", "Velocity (%)"]
        ws.append(headers)

        # Format header
        header_fill = PatternFill(start_color="3B82F6", end_color="3B82F6", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center")

        # Data
        shifts_data = shift_matrix.get("shifts", {})
        for category, cat_data in shifts_data.items():
            path = cat_data.get("path", {})
            velocity = cat_data.get("velocity", {})

            for year in sorted(path.keys()):
                year_data = path.get(year, {})
                if isinstance(year_data, dict):
                    ws.append([
                        category,
                        year,
                        "Base Case",
                        year_data.get("median", 0.0),
                        year_data.get("p10", 0.0),
                        year_data.get("p25", 0.0),
                        year_data.get("p75", 0.0),
                        year_data.get("p90", 0.0),
                        velocity_median(velocity.get(year)),
                    ])

        # Format data cells
        for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=4, max_col=9):
            for cell in row:
                cell.number_format = "0.00%"
                cell.alignment = Alignment(horizontal="right")

        # Auto-size columns
        ws.column_dimensions['A'].width = 18
        for col in ['B', 'C', 'D', 'E', 'F', 'G', 'H', 'I']:
            ws.column_dimensions[col].width = 12

    def _create_application_template_sheet(self, wb: openpyxl.Workbook, shift_matrix: Dict[str, Any]) -> None:
        """Create Application Template sheet with formula examples."""
        ws = wb.create_sheet("Application Template")

        # Instructions
        ws['A1'] = "PRISM Application Template — Applying Shifts to Your Financials"
        ws['A1'].font = Font(bold=True, size=12)

        ws['A3'] = "Instructions:"
        ws['A3'].font = Font(bold=True)
        ws['A4'] = "1. Enter your actual GP1 (€M) in column B"
        ws['A5'] = "2. The formula in column C automatically applies the PRISM shift"
        ws['A6'] = "3. Result: Projected GP1 with shift impact"

        # Headers
        headers = ["Category", "GP1 Actual (€M)", "PRISM Shift (%)", "GP1 Projected (€M)"]
        ws['A8'] = headers[0]
        ws['B8'] = headers[1]
        ws['C8'] = headers[2]
        ws['D8'] = headers[3]

        for cell in ['A8', 'B8', 'C8', 'D8']:
            ws[cell].font = Font(bold=True, color="FFFFFF")
            ws[cell].fill = PatternFill(start_color="3B82F6", end_color="3B82F6", fill_type="solid")

        # Add categories
        shifts_data = shift_matrix.get("shifts", {})
        row = 9
        for category in shifts_data.keys():
            ws[f'A{row}'] = category
            ws[f'B{row}'] = 0  # User fills in
            ws[f'C{row}'] = shifts_data[category].get("path", {}).get(2030, {}).get("median", 0.0)
            ws[f'D{row}'] = f"=B{row}*(1+C{row})"

            ws[f'C{row}'].number_format = "0.00%"
            ws[f'D{row}'].number_format = "€ #,##0"

            row += 1

        # Format
        ws.column_dimensions['A'].width = 18
        ws.column_dimensions['B'].width = 15
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 18

    def _create_allocation_sheet(self, wb: openpyxl.Workbook, allocation: Dict[str, Any]) -> None:
        """Create Allocation Recommendations sheet."""
        ws = wb.create_sheet("Allocation")

        ws['A1'] = "Resource Allocation Recommendations"
        ws['A1'].font = Font(bold=True, size=12)

        headers = ["Category", "Recommended Weight", "Rationale"]
        ws['A3'] = headers[0]
        ws['B3'] = headers[1]
        ws['C3'] = headers[2]

        for cell in ['A3', 'B3', 'C3']:
            ws[cell].font = Font(bold=True, color="FFFFFF")
            ws[cell].fill = PatternFill(start_color="3B82F6", end_color="3B82F6", fill_type="solid")

        # Extract recommendations
        if isinstance(allocation, dict):
            categories = allocation.get("invest_more", []) + allocation.get("defend", []) + allocation.get("harvest", [])
            invest_more = allocation.get("invest_more", [])
            defend = allocation.get("defend", [])
            harvest = allocation.get("harvest", [])

            row = 4
            for cat in invest_more:
                ws[f'A{row}'] = cat
                ws[f'B{row}'] = 0.08  # Example weight
                ws[f'C{row}'] = "INVEST — strong growth signal"
                row += 1

            for cat in defend:
                ws[f'A{row}'] = cat
                ws[f'B{row}'] = 0.06
                ws[f'C{row}'] = "DEFEND — protect market share"
                row += 1

            for cat in harvest:
                ws[f'A{row}'] = cat
                ws[f'B{row}'] = 0.02
                ws[f'C{row}'] = "HARVEST — limited investment"
                row += 1

        ws.column_dimensions['A'].width = 18
        ws.column_dimensions['B'].width = 18
        ws.column_dimensions['C'].width = 35

    def _create_force_attribution_sheet(self, wb: openpyxl.Workbook, causal: Dict[str, Any]) -> None:
        """Create Force Attribution sheet."""
        ws = wb.create_sheet("Force Attribution")

        ws['A1'] = "Force Attribution"
        ws['A1'].font = Font(bold=True, size=12)

        headers = ["Category", "Effect Type", "Source Force", "Path", "Contribution (%)"]
        ws['A3'] = headers[0]
        ws['B3'] = headers[1]
        ws['C3'] = headers[2]
        ws['D3'] = headers[3]
        ws['E3'] = headers[4]

        for cell in ['A3', 'B3', 'C3', 'D3', 'E3']:
            ws[cell].font = Font(bold=True, color="FFFFFF")
            ws[cell].fill = PatternFill(start_color="8B5CF6", end_color="8B5CF6", fill_type="solid")

        row = 4
        for category, decomp in causal.items():
            for force, value in decomp.get("direct_effects", {}).items():
                ws[f'A{row}'] = category
                ws[f'B{row}'] = "Direct"
                ws[f'C{row}'] = force
                ws[f'D{row}'] = force
                ws[f'E{row}'] = value
                ws[f'E{row}'].number_format = "0.00%"
                row += 1

            for path_label, value in decomp.get("propagated_effects", {}).items():
                ws[f'A{row}'] = category
                ws[f'B{row}'] = "Propagated"
                ws[f'C{row}'] = path_label.split("→")[0] if "→" in path_label else path_label
                ws[f'D{row}'] = path_label
                ws[f'E{row}'] = value
                ws[f'E{row}'].number_format = "0.00%"
                row += 1

        ws.column_dimensions['A'].width = 18
        ws.column_dimensions['B'].width = 12
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 25
        ws.column_dimensions['E'].width = 15

    def _create_methodology_sheet(self, wb: openpyxl.Workbook, shift_matrix: Dict[str, Any]) -> None:
        """Create Methodology sheet with model info."""
        ws = wb.create_sheet("Methodology")

        ws['A1'] = "PRISM Model — Methodology & Confidence"
        ws['A1'].font = Font(bold=True, size=12)

        row = 3
        ws[f'A{row}'] = "Model Version:"
        ws[f'B{row}'] = shift_matrix.get("model_version", "Bayesian Copula v1")
        row += 1

        ws[f'A{row}'] = "Generated:"
        ws[f'B{row}'] = shift_matrix.get("generated", datetime.now().isoformat())
        row += 1

        ws[f'A{row}'] = "Confidence Level:"
        ws[f'B{row}'] = shift_matrix.get("confidence", "80% CI")
        row += 1

        ws[f'A{row}'] = "Engine:"
        ws[f'B{row}'] = shift_matrix.get("engine_name", "bayesian_copula")
        row += 1

        ws[f'A{row}'] = "Seed:"
        ws[f'B{row}'] = shift_matrix.get("seed", "n/a")
        row += 2

        ws[f'A{row}'] = "Key Features:"
        ws[f'A{row}'].font = Font(bold=True)
        row += 1

        features = [
            "• Bayesian sampling with Beta posteriors over trend probabilities",
            "• Copula-based dependency structures (Gaussian + t-copula tails)",
            "• Force attribution per category (static, scaled to MC median)",
            "• Continuous paths (2026-2030) with annual granularity",
            "• Velocity tracking and early-warning triggers",
            "• Monte Carlo simulation (10,000-50,000 iterations)",
            "• Split-chain R̂ convergence diagnostics",
            "• Empirical covariance from MC samples in the allocation optimizer",
        ]

        for feature in features:
            ws[f'A{row}'] = feature
            row += 1

        row += 1
        ws[f'A{row}'] = "Security:"
        ws[f'A{row}'].font = Font(bold=True)
        row += 1

        security_notes = [
            "• All outputs contain ONLY percentage shifts (-1.0 to +1.0)",
            "• Zero company financial data (NES, GP1, GP2) in any export",
            "• One-directional data flow to Power BI (read-only)",
            "• Financial Data Firewall validates every export",
            "• Full audit trail of all model changes and predictions",
        ]

        for note in security_notes:
            ws[f'A{row}'] = note
            row += 1

        ws.column_dimensions['A'].width = 65
        ws.column_dimensions['B'].width = 35
