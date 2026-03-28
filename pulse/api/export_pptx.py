"""PowerPoint export for PULSE War Room results.

Generates professional executive decks with shift matrices, trend analysis,
and strategic recommendations. All values are percentages only.
"""

import logging
from datetime import datetime
from typing import Dict, List, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pulse.config import CATEGORIES, FORCES

logger = logging.getLogger(__name__)

# Design tokens matching React War Room
COLORS = {
    "primary": RGBColor(0, 113, 227),          # #0071E3
    "accent": RGBColor(123, 97, 255),          # #7B61FF (purple)
    "success": RGBColor(48, 209, 88),          # #30D158 (green)
    "danger": RGBColor(255, 69, 58),           # #FF453A (red)
    "warning": RGBColor(255, 159, 10),         # #FF9F0A (amber)
    "text": RGBColor(29, 29, 31),              # #1D1D1F
    "text_secondary": RGBColor(110, 110, 115), # #6E6E73
    "bg": RGBColor(255, 255, 255),             # white
    "bg_dark": RGBColor(245, 245, 247),        # #F5F5F7
    "bg_darker": RGBColor(239, 239, 239),      # #EFEFEF
}

FORCE_COLORS = {
    "Consumer": RGBColor(0, 113, 227),         # Blue
    "Customer": RGBColor(123, 97, 255),        # Purple
    "Technology": RGBColor(0, 180, 216),       # Cyan
    "Government": RGBColor(255, 159, 10),      # Amber
    "Environmental": RGBColor(48, 209, 88),    # Green
    "Competitive": RGBColor(255, 69, 58),      # Red
}


class PowerPointExporter:
    """Generate PULSE PowerPoint presentation."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def export(
        self,
        output_path: str,
        scenario: str,
        shifts: Dict[str, Any],
        trends: List[Dict[str, Any]],
        convergence: Optional[Dict[str, Any]] = None,
        allocation: Optional[Dict[str, Any]] = None,
        model_version: str = "v3.0",
        model_accuracy: float = 0.73,
    ) -> str:
        """
        Generate complete PowerPoint presentation.

        Args:
            output_path: Path to write .pptx file
            scenario: Scenario name (e.g., "Base Case", "Green Squeeze")
            shifts: Shift matrix {category: {year: {percentile: value}}}
            trends: List of trend objects with impact/probability
            convergence: Monte Carlo convergence diagnostics
            allocation: Allocation recommendations {category: weight}
            model_version: Model version string
            model_accuracy: Backtesting accuracy (0-1)

        Returns:
            Path to generated file
        """
        logger.info(f"Generating PowerPoint presentation: {scenario}")

        # Build slides
        self._add_title_slide(scenario, model_version)
        self._add_executive_summary_slide(shifts, trends, model_accuracy)
        self._add_shift_heatmap_slide(shifts)
        self._add_top_trends_slide(trends)
        self._add_strategic_implications_slide(shifts, allocation)

        # Save
        self.prs.save(output_path)
        logger.info(f"PowerPoint saved to {output_path}")
        return output_path

    def _add_title_slide(self, scenario: str, model_version: str):
        """Slide 1: Title slide with date and scenario."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["primary"]

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "PULSE War Room"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLORS["bg"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1.0))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = f"Profit Pool Shift Analysis — {scenario}"
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS["bg"]
        p.alignment = PP_ALIGN.CENTER

        # Metadata (date, version)
        meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.0))
        meta_frame = meta_box.text_frame
        meta_frame.word_wrap = True
        today = datetime.now().strftime("%B %d, %Y")
        p = meta_frame.paragraphs[0]
        p.text = f"{today} • Model {model_version}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(
        self, shifts: Dict[str, Any], trends: List[Dict[str, Any]], model_accuracy: float
    ):
        """Slide 2: Executive summary with key metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["bg"]

        # Header
        self._add_header(slide, "Executive Summary", 0.5)

        # Calculate key metrics
        median_shifts = []
        for cat, data in shifts.items():
            if not isinstance(data, dict):
                continue
            # Support both flat {year: {median}} and nested {path: {year: {median}}}
            year_data = data.get("path", data) if "path" in data else data
            if isinstance(year_data, dict) and "2030" in year_data and isinstance(year_data["2030"], dict):
                median = year_data["2030"].get("median", 0)
                if median is not None:
                    median_shifts.append((cat, median))

        if median_shifts:
            median_shifts.sort(key=lambda x: x[1], reverse=True)
            top_expansion = median_shifts[0] if median_shifts[0][1] > 0 else None
            top_contraction = [x for x in median_shifts if x[1] < 0][-1] if any(
                x[1] < 0 for x in median_shifts
            ) else None
            net_shift = sum(x[1] for x in median_shifts) / len(median_shifts) if median_shifts else 0
        else:
            top_expansion = None
            top_contraction = None
            net_shift = 0

        # Metrics grid (2x2)
        y_start = 1.2
        metrics = [
            ("Net Pool Shift", f"{net_shift:+.1%}", COLORS["primary"]),
            ("Model Accuracy", f"{model_accuracy:.0%}", COLORS["success"]),
            (
                "Top Expansion",
                f"{top_expansion[0]}: {top_expansion[1]:+.1%}" if top_expansion else "—",
                COLORS["success"],
            ),
            (
                "Top Contraction",
                f"{top_contraction[0]}: {top_contraction[1]:+.1%}" if top_contraction else "—",
                COLORS["danger"],
            ),
        ]

        for i, (label, value, color) in enumerate(metrics):
            x = 0.5 + (i % 2) * 4.75
            y = y_start + (i // 2) * 1.8
            self._add_metric_box(slide, x, y, label, value, color)

        # Key insights
        insights_y = 5.2
        insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(insights_y), Inches(9), Inches(1.8))
        insights_frame = insights_box.text_frame
        insights_frame.word_wrap = True

        p = insights_frame.paragraphs[0]
        p.text = "Key Insights"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]

        insights = []
        if top_expansion:
            insights.append(f"• Highest growth opportunity: {top_expansion[0]} ({top_expansion[1]:+.1%})")
        if top_contraction:
            insights.append(f"• Largest headwind: {top_contraction[0]} ({top_contraction[1]:+.1%})")
        if model_accuracy >= 0.7:
            insights.append(f"• Model shows strong backtested accuracy ({model_accuracy:.0%})")

        for insight in insights[:2]:
            p = insights_frame.add_paragraph()
            p.text = insight
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["text_secondary"]
            p.space_before = Pt(4)

    def _add_shift_heatmap_slide(self, shifts: Dict[str, Any]):
        """Slide 3: Category × Year heatmap table."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["bg"]

        self._add_header(slide, "Category Shift Paths (Median %)", 0.5)

        # Build table: categories × years
        years = [2026, 2027, 2028, 2029, 2030]
        rows = len(CATEGORIES) + 1  # +1 for header
        cols = len(years) + 1  # +1 for category names

        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(2.0)
        for i in range(1, cols):
            table.columns[i].width = Inches(1.2)

        # Header row
        for j, year in enumerate(years):
            cell = table.cell(0, j + 1)
            cell.text = str(year)
            self._format_header_cell(cell)

        # Category rows
        for i, cat_def in enumerate(CATEGORIES):
            row = i + 1
            cat_name = cat_def.name if hasattr(cat_def, "name") else (
                cat_def.get("name") if isinstance(cat_def, dict) else str(cat_def)
            )

            # Category name
            cell = table.cell(row, 0)
            cell.text = cat_name
            self._format_category_cell(cell)

            # Year values
            cat_key = cat_name
            for j, year in enumerate(years):
                cell = table.cell(row, j + 1)
                shift_data = shifts.get(cat_key, {})
                # Support both flat {year: {median}} and nested {path: {year: {median}}}
                path_data = shift_data.get("path", shift_data) if isinstance(shift_data, dict) and "path" in shift_data else shift_data
                year_data = path_data.get(str(year), {}) if isinstance(path_data, dict) else {}
                median = year_data.get("median", 0) if isinstance(year_data, dict) else 0

                if median is not None:
                    cell.text = f"{median:+.1%}"
                    self._format_shift_cell(cell, median)
                else:
                    cell.text = "—"
                    cell.text_frame.paragraphs[0].font.color.rgb = COLORS["text_secondary"]

    def _add_top_trends_slide(self, trends: List[Dict[str, Any]]):
        """Slide 4: Top 5 trends by impact × probability."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["bg"]

        self._add_header(slide, "Top Drivers: Trends by Impact", 0.5)

        # Score and rank trends
        scored_trends = []
        for trend in trends:
            impact = trend.get("impact", 3)
            probability = trend.get("probability", 3)
            score = (impact * probability) / 25.0  # Normalize to 0-1
            scored_trends.append((trend, score))

        scored_trends.sort(key=lambda x: x[1], reverse=True)
        top_trends = scored_trends[:5]

        # Table: 5 top trends
        rows = len(top_trends) + 1
        cols = 4

        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(1.5)

        # Header
        headers = ["Trend", "Force", "Direction", "Score"]
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            self._format_header_cell(cell)

        # Data rows
        for i, (trend, score) in enumerate(top_trends):
            row = i + 1
            name = trend.get("name", "Unknown")
            force = trend.get("force", "—")
            direction = trend.get("direction", "—")

            table.cell(row, 0).text = name
            table.cell(row, 1).text = force
            table.cell(row, 2).text = direction
            table.cell(row, 3).text = f"{score:.0%}"

            # Format force cell
            force_color = FORCE_COLORS.get(force, COLORS["text"])
            for col in range(cols):
                cell = table.cell(row, col)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                if col == 1:
                    cell.text_frame.paragraphs[0].font.color.rgb = force_color
                    cell.text_frame.paragraphs[0].font.bold = True

    def _add_strategic_implications_slide(
        self, shifts: Dict[str, Any], allocation: Optional[Dict[str, Any]] = None
    ):
        """Slide 5: Strategic implications and recommendations."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["bg"]

        self._add_header(slide, "Strategic Implications & Recommendations", 0.5)

        # Identify defend/harvest/invest categories
        defend = []
        harvest = []
        invest = []

        for cat, data in shifts.items():
            if not isinstance(data, dict):
                continue
            # Support both flat {year: {median}} and nested {path: {year: {median}}}
            year_data_container = data.get("path", data) if "path" in data else data
            if isinstance(year_data_container, dict) and "2030" in year_data_container and isinstance(year_data_container["2030"], dict):
                median = year_data_container["2030"].get("median", 0)
                if median is not None:
                    if -0.05 <= median <= 0.05:
                        defend.append((cat, median))
                    elif median < -0.05:
                        harvest.append((cat, median))
                    else:
                        invest.append((cat, median))

        # Bullet points
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = box.text_frame
        text_frame.word_wrap = True

        # Invest section
        p = text_frame.paragraphs[0]
        p.text = "Invest in Growth Categories"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["success"]

        for cat, shift in invest[:3]:
            p = text_frame.add_paragraph()
            p.text = f"• {cat}: +{shift:.1%} pool shift by 2030"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_before = Pt(4)

        # Defend section
        p = text_frame.add_paragraph()
        p.text = "Defend Core Business"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.space_before = Pt(12)

        for cat, _ in defend[:3]:
            p = text_frame.add_paragraph()
            p.text = f"• {cat}: Stable — maintain competitive position"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_before = Pt(4)

        # Harvest section
        p = text_frame.add_paragraph()
        p.text = "Harvest Declining Categories"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["danger"]
        p.space_before = Pt(12)

        for cat, shift in harvest[:3]:
            p = text_frame.add_paragraph()
            p.text = f"• {cat}: {shift:.1%} shift — optimize for cash flow"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["text"]
            p.level = 0
            p.space_before = Pt(4)

    # ─── Helper methods ────────────────────────────────────────

    def _add_header(self, slide, title: str, top_inches: float):
        """Add a standard slide header."""
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0),
            Inches(top_inches),
            Inches(10),
            Inches(0.6),
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLORS["bg_dark"]
        header_shape.line.color.rgb = COLORS["text_secondary"]

        text_frame = header_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.LEFT

    def _add_metric_box(self, slide, x: float, y: float, label: str, value: str, color: RGBColor):
        """Add a metric box (label + value)."""
        # Background
        box_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x),
            Inches(y),
            Inches(4.5),
            Inches(1.5),
        )
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = COLORS["bg_dark"]
        box_shape.line.color.rgb = color
        box_shape.line.width = Pt(2)

        # Label
        label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(4.1), Inches(0.5))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["text_secondary"]

        # Value
        value_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(4.1), Inches(0.7))
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color

    def _format_header_cell(self, cell):
        """Format table header cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS["bg"]
            paragraph.alignment = PP_ALIGN.CENTER

    def _format_category_cell(self, cell):
        """Format table category cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["bg_dark"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS["text"]

    def _format_shift_cell(self, cell, shift: float):
        """Format shift cell with color coding."""
        if shift > 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(220, 252, 231)  # Light green
            color = COLORS["success"]
        elif shift < 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(254, 226, 226)  # Light red
            color = COLORS["danger"]
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["bg_darker"]
            color = COLORS["text_secondary"]

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = color
            paragraph.alignment = PP_ALIGN.CENTER
